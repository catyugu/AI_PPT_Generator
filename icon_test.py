import os
import io
import cairosvg
from pptx import Presentation
from pptx.util import Inches


# --- 这是我们要测试的核心函数 ---
def convert_svg_to_png_stream(svg_path: str, color: str) -> io.BytesIO:
    """
    读取SVG文件，修改其描边颜色，然后将其转换为PNG格式的内存字节流。

    参数:
    svg_path (str): SVG文件的路径。
    color (str): 十六进制颜色代码，不带'#' (例如: 'EA4335')。

    返回:
    io.BytesIO: 包含PNG数据的内存字节流，如果失败则返回None。
    """
    print(f"开始处理SVG: {svg_path}")
    if not os.path.exists(svg_path):
        print(f"错误: 文件不存在 -> {svg_path}")
        return None

    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            svg_content = f.read()

        # Feather图标的默认颜色属性是 'currentColor'
        # 我们将其替换为指定的颜色，这让图标能融入任意PPT主题
        modified_svg_content = svg_content.replace('stroke="currentColor"', f'stroke="#{color}"')
        print(f"已将SVG颜色替换为: #{color}")

        # 创建一个内存中的字节流对象
        png_stream = io.BytesIO()

        # 使用cairosvg将修改后的SVG字符串转换为PNG，并写入内存流
        # output_width/height可以控制输出PNG的分辨率，更高分辨率更清晰
        cairosvg.svg2png(
            bytestring=modified_svg_content.encode('utf-8'),
            write_to=png_stream,
            output_width=256,
            output_height=256
        )

        # 将流的指针移回开头，以便pptx库可以读取它
        png_stream.seek(0)
        print("SVG成功转换为PNG内存流。")
        return png_stream
    except Exception as e:
        print(f"错误: 转换SVG时发生异常: {e}")
        return None


# --- 这是我们的测试主程序 ---
def run_test():
    """
    执行一个完整的测试：创建一个PPT，将一个变色的SVG图标插入其中。
    """
    print("--- 开始SVG插入测试 ---")

    # 1. 定义输入和输出
    svg_file = os.path.join("assets", "icons", "target.svg")
    output_ppt_file = "test_presentation.pptx"
    icon_color = "EA4335"  # 一个鲜艳的红色 (Google Red)

    # 2. 调用核心函数，转换SVG
    image_stream = convert_svg_to_png_stream(svg_file, icon_color)

    if not image_stream:
        print("测试失败：未能生成图片流。")
        return

    # 3. 创建一个新的PPT文稿
    prs = Presentation()
    # 使用一个空白页布局
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    print(f"已创建PPT '{output_ppt_file}' 并添加新幻灯片。")

    # 4. 将图片流添加到幻灯片中
    # 将图标放在幻灯片中央附近
    left = Inches(4)
    top = Inches(3)
    # 插入时可以只指定宽度，pptx会自动计算高度以保持纵横比
    width = Inches(2.0)

    try:
        pic = slide.shapes.add_picture(image_stream, left, top, width=width)
        print("成功将图片流添加到幻灯片！")
    except Exception as e:
        print(f"错误: 添加图片到幻灯片时失败: {e}")
    finally:
        # 确保关闭流
        image_stream.close()

    # 5. 保存PPT文件
    prs.save(output_ppt_file)
    print(f"--- 测试完成 ---")
    print(f"请检查文件: '{os.path.abspath(output_ppt_file)}'")


if __name__ == "__main__":
    run_test()