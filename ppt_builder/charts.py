# ppt_builder/charts.py
import logging
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def add_bar_chart(slide, chart_data, style):
    """Adds a bar chart to the slide based on provided data."""
    try:
        categories = chart_data.get("categories", ["C1", "C2", "C3"])
        series_name = chart_data.get("series_name", "Series 1")
        values = chart_data.get("values", [10, 20, 15])

        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = categories
        chart_data_obj.add_series(series_name, tuple(values))  # Values must be a tuple

        x, y, cx, cy = Inches(2), Inches(2.5), Inches(12), Inches(5.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
        )
        chart = graphic_frame.chart

        # Style the chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.plots[0].vary_by_categories = True
        chart.category_axis.tick_labels.font.size = Pt(14)
        chart.value_axis.tick_labels.font.size = Pt(14)
        chart.value_axis.has_major_gridlines = False

    except Exception as e:
        logging.error(f"Could not create chart: {e}")
        tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
        tb.text = f"Error: Invalid data provided for chart.\nDetails: {e}"