import logging
import os
import re
import io
from io import BytesIO

from PIL import Image, ImageOps, ImageDraw
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# [新增] 导入底层XML操作所需的工具
from pptx.oxml.ns import qn

from config import ICON_DIR
from ppt_builder.styles import px_to_emu, hex_to_rgb, PresentationStyle
import cairosvg
from pptx.util import Inches

# 形状类型映射
SHAPE_TYPE_MAP = {
    'rectangle': MSO_SHAPE.RECTANGLE,
    'oval': MSO_SHAPE.OVAL,
    'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    'star': MSO_SHAPE.STAR_5_POINT,
    'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
}

# 对齐方式映射
ALIGNMENT_MAP = {
    'LEFT': PP_ALIGN.LEFT,
    'CENTER': PP_ALIGN.CENTER,
    'RIGHT': PP_ALIGN.RIGHT,
    'JUSTIFY': PP_ALIGN.JUSTIFY,
}


def find_icon_path(keyword: str) -> str | None:
    """根据关键词在图标库中查找对应的SVG文件路径。"""
    # 简单的中英文关键词映射 (可选，但建议保留以增加AI的灵活性)
    keyword_map = {
        "创新": "cpu",
        "想法": "cpu",
        "增长": "trending-up",
        "目标": "target",
        "数据": "database",
        "警告": "alert-triangle",
        "设置": "settings",
        "链接": "link",
        "主页": "home"
    }

    normalized_keyword = keyword_map.get(keyword.lower(), keyword.lower())
    icon_path = os.path.join(ICON_DIR, f"{normalized_keyword}.svg")

    if os.path.exists(icon_path):
        return icon_path

    print(f"警告: 无法为关键词 '{keyword}' (解析为 '{normalized_keyword}.svg') 找到图标。将跳过。")
    return None


def _crop_to_circle(image_path: str):
    """
    使用Pillow将指定路径的图片处理成圆形，并返回一个内存中的图片流。
    """
    try:
        with Image.open(image_path) as img:
            img = img.convert("RGBA")
            size = (min(img.size), min(img.size))
            mask = Image.new('L', size, 0)
            draw = ImageDraw.Draw(mask)
            draw.ellipse((0, 0) + size, fill=255)
            output = ImageOps.fit(img, mask.size, centering=(0.5, 0.5))
            output.putalpha(mask)
            buffer = io.BytesIO()
            output.save(buffer, format='PNG')
            buffer.seek(0)
            return buffer
    except Exception as e:
        logging.error(f"处理图片为圆形时失败: {image_path} - {e}", exc_info=True)
        return None


def convert_svg_to_png_stream(svg_path: str, color: str) -> BytesIO | None:
    """读取SVG，用主题色替换其颜色，并转换为PNG内存流。"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            svg_content = f.read()

        # 将 'currentColor' 替换为指定的主题色
        modified_svg_content = svg_content.replace('stroke="currentColor"', f'stroke="#{color}"')

        png_stream = io.BytesIO()
        cairosvg.svg2png(
            bytestring=modified_svg_content.encode('utf-8'),
            write_to=png_stream,
            output_width=256  # 输出一个清晰的分辨率
        )
        png_stream.seek(0)
        return png_stream
    except Exception as e:
        print(f"错误: 转换SVG '{os.path.basename(svg_path)}' 时发生异常: {e}")
        return None


# ==================== [新增] 底层XML透明度设置函数 ====================
def _apply_transparency_to_color_format(color_format, opacity: float):
    """
    [Definitive Fix] Applies transparency by directly manipulating the underlying XML.
    This works for both solid and gradient fills.

    :param color_format: A ColorFormat object (e.g., fill.fore_color, stop.color).
    :param opacity: Opacity value from 0.0 (transparent) to 1.0 (opaque).
    """
    if not isinstance(opacity, (float, int)) or not (0.0 <= opacity <= 1.0):
        return

    # Convert opacity (1=opaque) to OpenXML alpha value (0=opaque, 100000=transparent)
    alpha_val = int((1.0 - opacity) * 100000)

    # The correct internal attribute to get the <a:srgbClr>, <a:schemeClr>, etc. element is `_color`
    color_element = color_format._color

    # Find the actual color definition element (e.g., srgbClr)
    color_type_element = color_element._srgbClr

    if color_type_element is None:
        color_type_element = color_element.schemeClr

    if color_type_element is None:
        logging.warning("Could not determine color type to apply transparency.")
        return

    # Find or create the <a:alpha> child element
    alpha_element = color_type_element.find(qn('a:alpha'))
    if alpha_element is None:
        alpha_element = color_type_element.makeelement(qn('a:alpha'))
        color_type_element.append(alpha_element)

    # Set the 'val' attribute to our calculated alpha value
    alpha_element.set('val', str(alpha_val))


def add_shape(slide, element_data: dict, style_manager: PresentationStyle):
    """
    [终极修复版] 添加形状并应用样式。
    - 使用底层XML操作，确保透明度对纯色和渐变填充都正确生效。
    """
    try:
        x, y, width, height = map(px_to_emu, [
            element_data.get('x', 50), element_data.get('y', 50),
            element_data.get('width', 200), element_data.get('height', 200)
        ])
        shape_type_str = element_data.get('shape_type', 'rectangle').lower()
        shape_type = SHAPE_TYPE_MAP.get(shape_type_str, MSO_SHAPE.RECTANGLE)
        style = element_data.get('style', {})

        shape = slide.shapes.add_shape(shape_type, x, y, width, height)
        fill = shape.fill
        line = shape.line

        # 统一获取透明度值
        opacity = style.get('opacity')
        has_opacity = isinstance(opacity, (float, int)) and 0 <= opacity <= 1

        # --- 设置填充 ---
        if 'gradient' in style:
            grad_info = style['gradient']
            fill.gradient()
            if 'angle' in grad_info:
                fill.gradient_angle = grad_info['angle']

            for i, hex_color in enumerate(grad_info.get('colors', [])):
                if i < len(fill.gradient_stops):
                    stop = fill.gradient_stops[i]
                    try:
                        stop.color.rgb = hex_to_rgb(hex_color)
                        if has_opacity:
                            _apply_transparency_to_color_format(stop.color, opacity)
                    except Exception as e:
                        logging.warning(f"设置渐变色标时出错: {e}")

            log_msg = "为形状应用了渐变填充"
            if has_opacity: log_msg += f" (透明度 opacity={opacity})"
            logging.info(log_msg)

        elif 'fill_color' in style and style['fill_color'] is not None:
            fill.solid()
            try:
                fill.fore_color.rgb = hex_to_rgb(style['fill_color'])
                if has_opacity:
                    _apply_transparency_to_color_format(fill.fore_color, opacity)
                    logging.info(f"为形状应用了纯色填充和透明度 (opacity={opacity})")
            except Exception as e:
                logging.warning(f"设置纯色填充时出错: {e}")
        else:
            fill.background()

        # --- 设置边框 ---
        if border_style := style.get('border'):
            try:
                line_color_hex = border_style.get('color', '#000000')
                line.color.rgb = hex_to_rgb(line_color_hex)
                line.width = Pt(border_style.get('width', 1))
            except (ValueError, IndexError):
                logging.warning(f"形状边框中提供了无效的十六进制颜色 '{line_color_hex}'。")
        else:
            line.fill.background()

        logging.info(f"添加 {shape_type_str} 形状完成。")
    except Exception as e:
        logging.error(f"添加形状时发生意外错误: {e}", exc_info=True)


# --- 其他函数 (add_text_box, add_image, add_chart, add_table) 保持不变 ---
# ... (将您文件中其他未改动的函数复制到这里)
def add_text_box(slide, element_data: dict, style_manager: PresentationStyle):
    """
    [已更新] 添加文本框，实现灵活的字体控制和项目符号列表。
    - 优先使用元素自身指定的字体 (`font.name`)。
    - 如果未指定，则根据 `font.type` ('heading'/'body') 回退到全局默认字体。
    - 兼容处理Markdown风格的加粗。
    - **[新功能]** 如果`content`是列表，则自动生成项目符号列表。
    """
    try:
        x, y, width, height = map(px_to_emu, [
            element_data.get('x', 50), element_data.get('y', 50),
            element_data.get('width', 1180), element_data.get('height', 100)
        ])

        txBox = slide.shapes.add_textbox(x, y, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.clear()

        content = element_data.get('content', '')
        style = element_data.get('style', {})
        font_style = style.get('font', {})

        # --- 字体决策逻辑 ---
        font_name = font_style.get('name')
        if not font_name:
            font_type = font_style.get('type', 'body')
            font_name = style_manager.heading_font if font_type == 'heading' else style_manager.body_font

        default_font_color = hex_to_rgb(font_style['color']) if 'color' in font_style else style_manager.text_color
        font_size_pt = font_style.get('size', 18)
        is_italic_from_json = font_style.get('italic', False)
        is_bold_from_json = font_style.get('bold', False)

        def apply_run_style(run, text_part, is_markdown_bold=False):
            """辅助函数，用于应用样式到文本块"""
            run.text = text_part
            font = run.font
            font.name = font_name
            font.size = Pt(font_size_pt)
            font.italic = is_italic_from_json
            font.color.rgb = default_font_color
            font.bold = is_bold_from_json or is_markdown_bold

        # ================== 新增逻辑：项目符号列表 ==================
        if isinstance(content, list):
            logging.info(f"检测到项目符号列表，共 {len(content)} 项。")
            for i, item_text in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.level = 0  # 可根据需要设置缩进级别

                if '**' in item_text:
                    parts = re.split(r'\*\*(.*?)\*\*', item_text)
                    for j, part in enumerate(parts):
                        if not part: continue
                        apply_run_style(p.add_run(), part, is_markdown_bold=(j % 2 == 1))
                else:
                    apply_run_style(p.add_run(), item_text)

                # 为每个段落设置对齐方式
                if alignment_str := style.get('alignment'):
                    p.alignment = ALIGNMENT_MAP.get(alignment_str.upper(), PP_ALIGN.LEFT)
                else:
                    p.alignment = PP_ALIGN.LEFT
        # ================== 原有逻辑：处理单个字符串 ==================
        else:
            p = tf.paragraphs[0]
            if '**' in content:
                parts = re.split(r'\*\*(.*?)\*\*', content)
                for i, part in enumerate(parts):
                    if not part: continue
                    apply_run_style(p.add_run(), part, is_markdown_bold=(i % 2 == 1))
            else:
                apply_run_style(p.add_run(), content)

            # 设置段落对齐
            if alignment_str := style.get('alignment'):
                p.alignment = ALIGNMENT_MAP.get(alignment_str.upper(), PP_ALIGN.LEFT)
            else:
                p.alignment = PP_ALIGN.LEFT

        log_content = str(content)
        logging.info(f"添加文本框 (字体: {font_name}): '{log_content[:30]}...'")
    except Exception as e:
        logging.error(f"添加文本框时出错: {e} | 原始元素数据: {element_data}", exc_info=True)


def add_image(slide, image_path: str, element_data: dict):
    """
    [最终版] 向幻灯片添加图片。
    - 如果 style.crop 为 'circle', 则将图片裁剪为圆形。
    - 否则，执行智能矩形裁剪以适应图框，避免拉伸。
    """
    try:
        if not image_path:
            logging.warning("图片路径为空，跳过添加图片。")
            return

        box_x, box_y, box_width, box_height = [
            element_data.get('x', 0), element_data.get('y', 0),
            element_data.get('width', 1280), element_data.get('height', 720)
        ]
        box_x_emu, box_y_emu, box_width_emu, box_height_emu = map(px_to_emu, [box_x, box_y, box_width, box_height])

        style = element_data.get('style', {})

        if style.get('crop') == 'circle':
            logging.info(f"检测到圆形裁剪请求，正在处理图片: {image_path}")
            circular_image_stream = _crop_to_circle(image_path)
            if circular_image_stream:
                diameter = min(box_width_emu, box_height_emu)
                slide.shapes.add_picture(circular_image_stream, box_x_emu, box_y_emu, width=diameter, height=diameter)
                logging.info(f"成功添加圆形图片: {image_path}")
            else:
                logging.error(f"圆形图片处理失败，无法添加图片: {image_path}")
        else:
            pic = slide.shapes.add_picture(image_path, box_x_emu, box_y_emu, width=box_width_emu, height=box_height_emu)
            with Image.open(image_path) as img:
                img_width_px, img_height_px = img.size
            if img_height_px == 0 or box_height == 0:
                logging.warning(f"图片或图框高度为零，无法计算宽高比，跳过裁剪: {image_path}")
                return
            img_aspect = img_width_px / img_height_px
            box_aspect = box_width / box_height
            if round(img_aspect, 2) != round(box_aspect, 2):
                if img_aspect > box_aspect:
                    crop_ratio = (1 - box_aspect / img_aspect) / 2
                    pic.crop_left = crop_ratio
                    pic.crop_right = crop_ratio
                else:
                    crop_ratio = (1 - img_aspect / box_aspect) / 2
                    pic.crop_top = crop_ratio
                    pic.crop_bottom = crop_ratio
            logging.info(f"从路径添加并智能裁剪矩形图片: {image_path}")

    except FileNotFoundError:
        logging.error(f"图片文件未找到: {image_path}")
    except Exception as e:
        logging.error(f"添加或裁剪图片 {image_path} 时出错: {e}", exc_info=True)


# 在 elements.py 文件中，请用下面的函数替换旧的 add_chart 函数

# 在 elements.py 文件中，请用下面的函数替换旧的 add_chart 函数

def add_chart(slide, element_data: dict, style_manager: PresentationStyle):
    """
    [终极美化版] 添加图表并进行深度样式化，以实现商业级报告外观。
    """
    try:
        x, y, width, height = map(px_to_emu, [element_data.get('x', 100), element_data.get('y', 150),
                                              element_data.get('width', 1080), element_data.get('height', 450)])
        chart_type_map = {'BAR': XL_CHART_TYPE.COLUMN_CLUSTERED, 'PIE': XL_CHART_TYPE.PIE, 'LINE': XL_CHART_TYPE.LINE}
        chart_type = chart_type_map.get(element_data.get('chart_type', 'bar').upper())

        chart_data = ChartData()
        chart_data_info = element_data.get('data', {})
        chart_data.categories = chart_data_info.get('categories', [])
        for series_data in chart_data_info.get('series', []):
            chart_data.add_series(series_data.get('name', ''), series_data.get('values', []))

        chart = slide.shapes.add_chart(chart_type, x, y, width, height, chart_data).chart

        # --- 1. 标题和图例 ---
        if title_text := element_data.get('title'):
            chart.has_title = True
            p = chart.chart_title.text_frame.paragraphs[0]
            p.font.size, p.font.bold = Pt(20), True
            p.font.color.rgb = style_manager.text_color
            p.font.name = style_manager.heading_font

        # [最终修复] 强制为图表创建图例
        chart.has_legend = True
        chart.legend.position, chart.legend.include_in_layout = XL_LEGEND_POSITION.BOTTOM, False
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = style_manager.text_color

        # --- 2. 绘图区和数据标签 ---
        plot = chart.plots[0]
        # 仅对非饼图设置 vary_by_categories = False
        if chart_type != XL_CHART_TYPE.PIE:
            plot.vary_by_categories = False

        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size, data_labels.font.bold = Pt(14), True
        data_labels.font.color.rgb = RGBColor(255, 255,
                                              255) if chart_type == XL_CHART_TYPE.PIE else style_manager.text_color
        if chart_type == XL_CHART_TYPE.PIE:
            data_labels.show_percentage = True
            data_labels.number_format = '0%'

        # --- 3. [核心] 系列的深度样式化 ---
        if chart_type in [XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE]:
            for i, series in enumerate(getattr(plot, 'series', [])):
                series_color = style_manager.get_chart_color(i)
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = series_color

                line = series.format.line
                if chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                    line.color.rgb = RGBColor(255, 255, 255)
                    line.width = Pt(0.75)
                elif chart_type == XL_CHART_TYPE.LINE:
                    line.color.rgb = series_color
                    line.width = Pt(2.5)
                    series.smooth = True
                    marker = series.marker
                    marker.style, marker.size = XL_MARKER_STYLE.CIRCLE, 8
                    marker.format.fill.solid()
                    marker.format.fill.fore_color.rgb = series_color
                    marker.format.line.color.rgb = RGBColor(255, 255, 255)
                    marker.format.line.width = Pt(1.0)

        elif chart_type == XL_CHART_TYPE.PIE and hasattr(plot, 'series') and plot.series:
            for j, point in enumerate(plot.series[0].points):
                point_color = style_manager.get_chart_color(j)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = point_color
                point.format.line.color.rgb = RGBColor(255, 255, 255)
                point.format.line.width = Pt(1.5)

        # --- 4. 坐标轴样式 ---
        if chart_type != XL_CHART_TYPE.PIE:
            for axis in [chart.category_axis, chart.value_axis]:
                axis.tick_labels.font.size = Pt(12)
                axis.tick_labels.font.color.rgb = style_manager.text_color
            if chart.value_axis.has_major_gridlines:
                chart.value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb("#E0E0E0")

        logging.info("成功添加并深度美化了图表。")
    except Exception as e:
        logging.error(f"添加图表时出错: {e}", exc_info=True)


def add_table(slide, element_data: dict, style_manager: PresentationStyle):
    """添加表格并应用样式。"""
    try:
        x, y, width, height = map(px_to_emu, [
            element_data.get('x', 100), element_data.get('y', 150),
            element_data.get('width', 1080), element_data.get('height', 420)
        ])

        headers = element_data.get('headers', [])
        rows_data = element_data.get('rows', [])
        if not headers or not rows_data:
            logging.warning("表格数据缺少表头或行数据，已跳过。")
            return

        num_rows, num_cols = len(rows_data) + 1, len(headers)
        shape = slide.shapes.add_table(num_rows, num_cols, x, y, width, height)
        table = shape.table

        for c in range(num_cols):
            table.columns[c].width = int(width / num_cols)
        for r in range(num_rows):
            table.rows[r].height = int(height / num_rows)

        style = element_data.get('style', {})
        header_color = hex_to_rgb(style.get('header_color')) if 'header_color' in style else style_manager.primary
        row_colors = [hex_to_rgb(c) for c in style.get('row_colors', [])]

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.name = style_manager.heading_font

        for r, row_data in enumerate(rows_data):
            for c, cell_data in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = str(cell_data)
                if row_colors:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_colors[r % len(row_colors)]
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = style_manager.text_color
                p.font.name = style_manager.body_font

        logging.info(f"添加了包含 {len(rows_data)} 行的表格。")
    except Exception as e:
        logging.error(f"添加表格时出错: {e}", exc_info=True)


def add_icon(slide, element_data, presentation_size, theme_colors):
    """
    在幻灯片上添加一个经过主题色渲染的图标。
    这个函数是外部调用的主入口。
    """
    keyword = element_data.get("icon_keyword")
    if not keyword:
        print("警告: 图标元素缺少 'icon_keyword'。将跳过。")
        return

    svg_path = find_icon_path(keyword)
    if not svg_path:
        return

    # 从PPT设计方案中获取一个强调色来渲染图标，确保风格统一
    # 您可以根据需要制定更复杂的逻辑，例如使用 accent2, accent3 等
    icon_color = theme_colors.get("accent1", "000000")  # 默认使用强调色1，如果没有则用黑色

    png_image_stream = convert_svg_to_png_stream(svg_path, icon_color)
    if not png_image_stream:
        return

    pres_width_emu, pres_height_emu = presentation_size
    pres_width_inches = pres_width_emu / 914400
    pres_height_inches = pres_height_emu / 914400

    left = Inches(element_data["x"] * pres_width_inches)
    top = Inches(element_data["y"] * pres_height_inches)
    width = Inches(element_data["width"] * pres_width_inches)

    try:
        slide.shapes.add_picture(png_image_stream, left, top, width=width)
        print(f"成功添加SVG图标: '{keyword}' (颜色: #{icon_color})")
    except Exception as e:
        print(f"错误: 从流添加图片 '{keyword}' 时发生异常: {e}")
    finally:
        png_image_stream.close()
