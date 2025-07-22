# ppt_builder/layouts.py
import logging
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from ppt_builder.shapes import add_image_with_fallback, add_overlay, add_standard_header
from ppt_builder.charts import add_bar_chart


def _get_sanitized_content_list(content):
    """Sanitizes AI content to always return a list of strings."""
    if not content:
        return []
    if isinstance(content, str):
        return [content]
    if isinstance(content, list):
        return [str(item) for item in content if item]
    return []


def _get_alignment_enum(align_str, default=PP_ALIGN.LEFT):
    """Converts a string to a PP_ALIGN enum."""
    return {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }.get(str(align_str).lower(), default)


def create_cover_slide(slide, page_data, style, prs):
    """[UPGRADED] Cover slide now supports variable opacity and text alignment."""
    layout_options = page_data.get('layout_options', {})

    keyword = page_data.get('image_keyword', 'abstract technology')
    add_image_with_fallback(slide, keyword, 0, 0, prs.slide_width, prs.slide_height)

    opacity = layout_options.get('overlay_opacity', 0.4)
    add_overlay(slide, prs.slide_width, prs.slide_height, style.primary, opacity)

    alignment = _get_alignment_enum(layout_options.get('text_alignment'), default=PP_ALIGN.LEFT)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
    p = title_box.text_frame.paragraphs[0]
    p.text = page_data.get('title', 'Presentation Title')
    p.font.name = style.font_heading
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = alignment

    content_list = _get_sanitized_content_list(page_data.get('content'))
    subtitle_text = content_list[0] if content_list else style.design_concept
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.75), Inches(14), Inches(1))
    p_sub = subtitle_box.text_frame.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = style.font_body
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = RGBColor(220, 220, 220)
    p_sub.alignment = alignment


def create_section_header_slide(slide, page_data, style, prs):
    """Creates a bold section divider with a large title and accent shape."""
    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(0.2), Inches(2))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = style.accent
    accent_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(12), Inches(2))
    p = title_box.text_frame.paragraphs[0]
    p.text = page_data.get('title', 'Section Header')
    p.font.name = style.font_heading
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = style.primary


def create_title_content_slide(slide, page_data, style, prs):
    """[UPGRADED] Title/Content slide now supports multiple columns."""
    add_standard_header(slide, page_data, style)
    layout_options = page_data.get('layout_options', {})

    content_list = _get_sanitized_content_list(page_data.get('content'))
    if not content_list:
        return

    num_columns = layout_options.get('columns', 1)

    if not isinstance(num_columns, int) or not (1 <= num_columns <= 3):
        num_columns = 1

    if num_columns == 1:
        body_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(6))
        tf = body_box.text_frame
        tf.word_wrap = True
        for point in content_list:
            p = tf.add_paragraph()
            p.text = point
            p.font.name = style.font_body
            p.font.size = Pt(22)
            p.font.color.rgb = style.text
            p.level = 0
            p.space_before = Pt(12)
    else:
        col_width = (prs.slide_width - Inches(2)) / num_columns
        points_per_col = (len(content_list) + num_columns - 1) // num_columns

        for i in range(num_columns):
            col_left = Inches(1) + i * col_width
            start_index = i * points_per_col
            end_index = start_index + points_per_col
            col_points = content_list[start_index:end_index]

            if not col_points: continue

            col_box = slide.shapes.add_textbox(col_left, Inches(2.5), col_width - Inches(0.5), Inches(6))
            tf = col_box.text_frame
            tf.word_wrap = True
            for point in col_points:
                p = tf.add_paragraph()
                p.text = point
                p.font.name = style.font_body
                p.font.size = Pt(20)
                p.font.color.rgb = style.text
                p.level = 0
                p.space_before = Pt(10)


def create_image_text_split_slide(slide, page_data, style, prs):
    """[UPGRADED] Split slide now supports variable ratios and image positions."""
    add_standard_header(slide, page_data, style)
    layout_options = page_data.get('layout_options', {})

    image_position = layout_options.get('image_position', 'right')
    ratio_str = layout_options.get('split_ratio', '50/50')

    try:
        r1, r2 = map(int, ratio_str.split('/'))
        total_ratio = r1 + r2
        if total_ratio == 0: raise ValueError()
    except (ValueError, IndexError, TypeError):
        r1, r2 = 50, 50
        total_ratio = 100

    content_width = prs.slide_width - Inches(2)
    gap = Inches(0.5)

    img_ratio_w = (content_width - gap) * (r1 / total_ratio)
    txt_ratio_w = (content_width - gap) * (r2 / total_ratio)

    if image_position == 'left':
        img_w, txt_w = img_ratio_w, txt_ratio_w
        img_left, txt_left = Inches(1), Inches(1) + img_w + gap
    else:
        txt_w, img_w = img_ratio_w, txt_ratio_w
        txt_left, img_left = Inches(1), Inches(1) + txt_w + gap

    add_image_with_fallback(slide, page_data.get('image_keyword', 'business'), img_left, Inches(2.2), img_w, Inches(6))

    text_box = slide.shapes.add_textbox(txt_left, Inches(2.2), txt_w, Inches(6))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    content_list = _get_sanitized_content_list(page_data.get('content'))
    for point in content_list:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = style.font_body
        p.font.size = Pt(20)
        p.font.color.rgb = style.text
        p.level = 0
        p.space_before = Pt(10)


def create_full_bleed_image_quote_slide(slide, page_data, style, prs):
    """Creates a slide with a full-screen image and a large quote."""
    add_image_with_fallback(slide, page_data.get('image_keyword', 'inspiration'), 0, 0, prs.slide_width,
                            prs.slide_height)
    add_overlay(slide, prs.slide_width, prs.slide_height, RGBColor(0, 0, 0), 0.6)

    quote_text = _get_sanitized_content_list(page_data.get('content', ['"Your quote here"']))[0]
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(13), Inches(3))
    p = quote_box.text_frame.paragraphs[0]
    p.text = quote_text
    p.font.name = style.font_heading
    p.font.size = Pt(44)
    p.font.italic = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def create_three_points_icons_slide(slide, page_data, style, prs):
    """Creates a slide highlighting three points with placeholder icons."""
    add_standard_header(slide, page_data, style)

    points = page_data.get('content', [])
    if not isinstance(points, list): points = []

    num_points = len(points)
    col_width = prs.slide_width / max(num_points, 1)

    for i, point_data in enumerate(points[:3]):
        left = (i * col_width)

        icon_keyword = point_data.get('icon_keyword', 'point') if isinstance(point_data, dict) else 'point'
        point_text = point_data.get('text', str(point_data)) if isinstance(point_data, dict) else str(point_data)

        icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + col_width / 2 - Inches(0.75), Inches(3), Inches(1.5),
                                            Inches(1.5))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = style.accent
        icon_shape.line.fill.background()

        desc_box = slide.shapes.add_textbox(left, Inches(5), col_width, Inches(3))
        p = desc_box.text_frame.paragraphs[0]
        p.text = point_text
        p.font.name = style.font_body
        p.font.size = Pt(18)
        p.font.color.rgb = style.text
        p.alignment = PP_ALIGN.CENTER


def create_bar_chart_slide(slide, page_data, style, prs):
    """Creates a slide with a bar chart."""
    add_standard_header(slide, page_data, style)
    chart_content = page_data.get('content')
    if isinstance(chart_content, dict):
        add_bar_chart(slide, chart_content, style)
    else:
        logging.warning("Bar chart content is not a valid dictionary.")
        tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
        tb.text = "错误：图表数据无效。"


def create_timeline_slide(slide, page_data, style, prs):
    """Creates a simple horizontal timeline."""
    add_standard_header(slide, page_data, style)
    events = page_data.get('content', [])
    if not isinstance(events, list) or not events: return

    num_events = len(events)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.5), Inches(14), Inches(0.05))
    line.fill.solid();
    line.fill.fore_color.rgb = style.secondary

    total_width = Inches(14)
    step_width = total_width / num_events

    for i, event in enumerate(events):
        left_pos = Inches(1) + (i * step_width) + (step_width / 2)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_pos - Inches(0.25), Inches(4.25), Inches(0.5), Inches(0.5))
        circle.fill.solid();
        circle.fill.fore_color.rgb = style.primary
        circle.line.fill.background()

        step_text = event.get('step', f"Step {i + 1}") if isinstance(event, dict) else f"Step {i + 1}"
        desc_text = event.get('description', '') if isinstance(event, dict) else str(event)

        tb = slide.shapes.add_textbox(left_pos - (step_width / 2), Inches(5.0), step_width, Inches(3))
        tf = tb.text_frame
        tf.word_wrap = True
        p_step = tf.paragraphs[0]
        p_step.text = step_text
        p_step.font.name = style.font_heading
        p_step.font.size = Pt(18);
        p_step.font.bold = True
        p_step.alignment = PP_ALIGN.CENTER

        p_desc = tf.add_paragraph()
        p_desc.text = desc_text
        p_desc.font.name = style.font_body
        p_desc.font.size = Pt(14)
        p_desc.alignment = PP_ALIGN.CENTER


def create_team_intro_slide(slide, page_data, style, prs):
    """[UPGRADED] Now handles both dict and simple string content for members."""
    add_standard_header(slide, page_data, style)
    members = page_data.get('content', [])
    if not isinstance(members, list): members = []

    num_members = len(members)
    if num_members == 0: return

    cols = min(num_members, 4)
    col_width = prs.slide_width / cols

    for i, member in enumerate(members[:cols]):
        left = (i * col_width)

        # [NEW] Check the format of 'member'
        if isinstance(member, dict):
            name = member.get('name', '姓名')
            title = member.get('title', '职位')
        else:  # Fallback for simple string content
            name = str(member)
            title = ''

        img_placeholder = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + col_width / 2 - Inches(1), Inches(2.5),
                                                 Inches(2), Inches(2))
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = style.accent
        img_placeholder.line.fill.background()

        name_box = slide.shapes.add_textbox(left, Inches(4.7), col_width, Inches(0.8))
        p_name = name_box.text_frame.paragraphs[0]
        p_name.text = name  # Safely use the parsed name
        p_name.font.name = style.font_heading
        p_name.font.size = Pt(24)
        p_name.font.bold = True
        p_name.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(left, Inches(5.5), col_width, Inches(1.5))
        p_title = title_box.text_frame.paragraphs[0]
        p_title.text = title  # Safely use the parsed title
        p_title.font.name = style.font_body
        p_title.font.size = Pt(16)
        p_title.alignment = PP_ALIGN.CENTER


def create_icon_grid_slide(slide, page_data, style, prs):
    """[UPGRADED] Now handles both dict and simple string content for grid items."""
    add_standard_header(slide, page_data, style)
    items = page_data.get('content', [])
    if not isinstance(items, list): items = []
    if not items: return

    cols, rows = 3, 2
    col_width, row_height = Inches(5), Inches(3)
    start_left, start_top = Inches(0.5), Inches(2.5)

    for i, item in enumerate(items[:cols * rows]):
        row = i // cols
        col = i % cols
        left = start_left + col * col_width
        top = start_top + row * row_height

        # [NEW] Check the format of 'item'
        if isinstance(item, dict):
            text = item.get('text', '项目说明')
        else:  # Fallback for simple string content
            text = str(item)

        icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(1.75), top, Inches(1.5),
                                            Inches(1.5))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = style.primary
        icon_shape.line.fill.background()

        text_box = slide.shapes.add_textbox(left, top + Inches(1.7), col_width, Inches(1))
        p = text_box.text_frame.paragraphs[0]
        p.text = text  # Safely use the parsed text
        p.font.name = style.font_body
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        text_box.text_frame.word_wrap = True

def create_process_flow_slide(slide, page_data, style, prs):
    """[UPGRADED] Now handles both dict and simple string content for process steps."""
    add_standard_header(slide, page_data, style)
    steps = page_data.get('content', [])
    if not isinstance(steps, list) or not steps: return

    num_steps = len(steps)
    total_width, gap_width = Inches(15), Inches(0.5)
    box_width = (total_width - (num_steps - 1) * gap_width) / num_steps
    box_height, top = Inches(3), Inches(4)

    for i, step_data in enumerate(steps[:5]):
        left = Inches(0.5) + i * (box_width + gap_width)

        # [NEW] Check the format of 'step_data'
        if isinstance(step_data, dict):
            step_text = step_data.get('step', f'步骤 {i + 1}')
            desc_text = step_data.get('description', '')
        else:  # Fallback for simple string content
            parts = str(step_data).split(':', 1)
            step_text = parts[0]
            desc_text = parts[1].strip() if len(parts) > 1 else ''

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = style.primary
        box.line.fill.background()

        tf = box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

        p_step = tf.paragraphs[0]
        p_step.text = step_text  # Safely use parsed text
        p_step.font.name = style.font_heading
        p_step.font.bold = True
        p_step.font.size = Pt(20)
        p_step.font.color.rgb = RGBColor(255, 255, 255)
        p_step.alignment = PP_ALIGN.CENTER

        p_desc = tf.add_paragraph()
        p_desc.text = desc_text  # Safely use parsed text
        p_desc.font.name = style.font_body
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = RGBColor(220, 220, 220)
        p_desc.alignment = PP_ALIGN.CENTER

        if i < num_steps - 1:
            arrow_left = left + box_width
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, top + box_height / 2 - Inches(0.25),
                                           gap_width, Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = style.accent
            arrow.line.fill.background()

def create_thank_you_slide(slide, page_data, style, prs):
    """Creates a dedicated 'Thank You' slide that doesn't require content."""
    keyword = page_data.get('image_keyword', 'abstract farewell')
    add_image_with_fallback(slide, keyword, 0, 0, prs.slide_width, prs.slide_height)
    add_overlay(slide, prs.slide_width, prs.slide_height, style.primary, 0.4)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(2))
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_box.text_frame.paragraphs[0]
    p.text = page_data.get('title', '谢谢观看')
    p.font.name = style.font_heading
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    content_list = _get_sanitized_content_list(page_data.get('content'))
    if content_list:
        subtitle_text = "\n".join(content_list)
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1.5))
        p_sub = subtitle_box.text_frame.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = style.font_body
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(220, 220, 220)
        p_sub.alignment = PP_ALIGN.CENTER