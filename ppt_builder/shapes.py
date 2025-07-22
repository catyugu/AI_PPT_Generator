# ppt_builder/shapes.py
import logging
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def add_image_with_fallback(slide, keyword, left, top, width, height):
    """Adds an image to the slide, handling potential fetch errors."""
    from image_service import get_image_stream # Defer import to avoid circular dependency
    image_stream = get_image_stream(keyword)
    if image_stream:
        try:
            slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        except Exception as e:
            logging.error(f"Could not add picture for keyword '{keyword}': {e}")


def add_overlay(slide, width, height, rgb_color, transparency):
    """Adds a semi-transparent overlay."""
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    fill = overlay.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color
    fill.transparency = transparency
    overlay.line.fill.background()


def add_standard_header(slide, page_data, style):
    """Adds a consistent title and decorative line to a slide."""
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.25))
    p = title_shape.text_frame.paragraphs[0]
    p.text = page_data.get('title', 'Slide Title')
    p.font.name = style.font_heading
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = style.primary

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.7), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = style.secondary
    line.line.fill.background()