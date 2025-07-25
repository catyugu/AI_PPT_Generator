import logging
import os

from pptx import Presentation
# 确保 elements 模块被正确导入
from ppt_builder import elements, animation
from ppt_builder.styles import PresentationStyle, px_to_emu

# --- 修改1：为icon元素添加图层顺序 ---
ELEMENT_DEFAULT_LAYER_ORDER = {
    'image': 10,  # 图片通常在背景之上，但可能在形状之下
    'shape': 20,  # 形状可以作为背景或蒙版，也可以在图片之上
    'icon': 30,   # 图标通常在形状和图片之上
    'chart': 40,  # 图表通常在图标之上
    'table': 40,  # 表格通常在图标之上
    'text_box': 50, # 文本框永远在最上层，确保可读性
    'default': 0  # 为未知类型提供默认层级 (最底层)
}


class SlideRenderer:
    """负责将单页幻灯片的数据渲染到演示文稿中。"""

    def __init__(self, prs: Presentation, style_manager: PresentationStyle, background_image_path: str | None):
        """
        初始化渲染器。
        :param prs: 演示文稿对象。
        :param style_manager: 全局样式管理器。
        :param background_image_path: 全局背景图片的路径，如果无则为None。
        """
        self.prs = prs
        self.style_manager = style_manager
        self.background_image_path = background_image_path
        logging.info("SlideRenderer已使用样式管理器和背景信息初始化。")

    def _add_background_image(self, slide, image_path: str):
        """
        [FIXED] Adds a background image and moves it to the back.
        Now includes checks to prevent errors with invalid paths.
        """
        if not image_path:
            # 静默处理，因为没有全局背景是正常情况
            return
        if not os.path.exists(image_path):
            logging.error(f"Background image not found at path: {image_path}. Skipping.")
            return

        try:
            picture = slide.shapes.add_picture(
                image_path, 0, 0,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
            # 将图片移动到最底层
            pic_element = picture.element
            slide.shapes._spTree.remove(pic_element)
            slide.shapes._spTree.insert(0, pic_element)

            logging.info(f"Successfully added background image: {image_path}")
        except Exception as e:
            logging.error(f"An unexpected error occurred while adding background image '{image_path}': {e}",
                          exc_info=True)

    def render_slide(self, slide_data: dict, image_service):
        """根据给定的数据渲染一张幻灯片。"""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # 首先添加全局背景图片（如果存在）
        self._add_background_image(slide, self.background_image_path)

        # ===================== 核心修改：元素排序 =====================
        elements_to_render = slide_data.get('elements', [])
        shapes_map = {}

        # 根据z_index排序，如果z_index不存在则回退到默认层级
        elements_to_render.sort(
            key=lambda e: e.get('z_index', ELEMENT_DEFAULT_LAYER_ORDER.get(e.get('type'), ELEMENT_DEFAULT_LAYER_ORDER['default']))
        )
        logging.info("页面元素已按图层顺序重排，渲染开始...")

        # 遍历排序后的列表进行渲染
        for element in elements_to_render:
            element_type = element.get('type')
            element_id = element.get('id')
            new_shape = None
            try:
                if element_type in ['text_box', 'text']:
                    new_shape = elements.add_text_box(slide, element, self.style_manager)

                elif element_type == 'image':
                    if image_keyword := element.get('image_keyword'):
                        opacity = element.get('style', {}).get('opacity', 1.0)
                        image_path = image_service.generate_image(image_keyword, opacity)
                        if image_path:
                            new_shape = elements.add_image(slide, image_path, element)
                        else:
                            logging.warning(f"无法为关键词生成图片: '{image_keyword}'。已跳过此元素。")
                    else:
                        logging.warning("图片元素缺少 'image_keyword'，已跳过。")

                elif element_type == 'shape':
                    new_shape = elements.add_shape(slide, element, self.style_manager)

                elif element_type == 'chart':
                    new_shape = elements.add_chart(slide, element, self.style_manager)

                elif element_type == 'table':
                    new_shape = elements.add_table(slide, element, self.style_manager)

                # --- 修改2：添加对icon类型的支持 ---
                elif element_type == 'icon':
                    new_shape = elements.add_icon(
                        slide,
                        element,
                        (self.prs.slide_width, self.prs.slide_height),
                        self.style_manager
                    )
                # -----------------------------------
                else:
                    logging.warning(f"不支持的元素类型: '{element_type}'。")

                if element_id and new_shape:
                    shapes_map[element_id] = new_shape

            except Exception as e:
                logging.error(f"渲染类型为 '{element_type}' 的元素失败: {e}", exc_info=True)
                # --- 第二步：按序列添加动画 ---
        if 'animation_sequence' in slide_data:
            for anim_step in slide_data['animation_sequence']:
                import logging
                import os

                from pptx import Presentation
                # 确保 elements 模块被正确导入
                from ppt_builder import elements, animation
                from ppt_builder.styles import PresentationStyle, px_to_emu

                # --- 修改1：为icon元素添加图层顺序 ---
                ELEMENT_DEFAULT_LAYER_ORDER = {
                    'image': 10,  # 图片通常在背景之上，但可能在形状之下
                    'shape': 20,  # 形状可以作为背景或蒙版，也可以在图片之上
                    'icon': 30,  # 图标通常在形状和图片之上
                    'chart': 40,  # 图表通常在图标之上
                    'table': 40,  # 表格通常在图标之上
                    'text_box': 50,  # 文本框永远在最上层，确保可读性
                    'default': 0  # 为未知类型提供默认层级 (最底层)
                }

                class SlideRenderer:
                    """负责将单页幻灯片的数据渲染到演示文稿中。"""

                    def __init__(self, prs: Presentation, style_manager: PresentationStyle,
                                 background_image_path: str | None):
                        """
                        初始化渲染器。
                        :param prs: 演示文稿对象。
                        :param style_manager: 全局样式管理器。
                        :param background_image_path: 全局背景图片的路径，如果无则为None。
                        """
                        self.prs = prs
                        self.style_manager = style_manager
                        self.background_image_path = background_image_path
                        logging.info("SlideRenderer已使用样式管理器和背景信息初始化。")

                    def _add_background_image(self, slide, image_path: str):
                        """
                        [FIXED] Adds a background image and moves it to the back.
                        Now includes checks to prevent errors with invalid paths.
                        """
                        if not image_path:
                            # 静默处理，因为没有全局背景是正常情况
                            return
                        if not os.path.exists(image_path):
                            logging.error(f"Background image not found at path: {image_path}. Skipping.")
                            return

                        try:
                            picture = slide.shapes.add_picture(
                                image_path, 0, 0,
                                width=self.prs.slide_width,
                                height=self.prs.slide_height
                            )
                            # 将图片移动到最底层
                            pic_element = picture.element
                            slide.shapes._spTree.remove(pic_element)
                            slide.shapes._spTree.insert(0, pic_element)

                            logging.info(f"Successfully added background image: {image_path}")
                        except Exception as e:
                            logging.error(
                                f"An unexpected error occurred while adding background image '{image_path}': {e}",
                                exc_info=True)

                    def render_slide(self, slide_data: dict, image_service):
                        """根据给定的数据渲染一张幻灯片。"""
                        blank_slide_layout = self.prs.slide_layouts[6]
                        slide = self.prs.slides.add_slide(blank_slide_layout)

                        # 首先添加全局背景图片（如果存在）
                        self._add_background_image(slide, self.background_image_path)

                        # ===================== 核心修改：元素排序 =====================
                        elements_to_render = slide_data.get('elements', [])
                        shapes_map = {}

                        # 根据z_index排序，如果z_index不存在则回退到默认层级
                        elements_to_render.sort(
                            key=lambda e: e.get('z_index', ELEMENT_DEFAULT_LAYER_ORDER.get(e.get('type'),
                                                                                           ELEMENT_DEFAULT_LAYER_ORDER[
                                                                                               'default']))
                        )
                        logging.info("页面元素已按图层顺序重排，渲染开始...")

                        # 遍历排序后的列表进行渲染
                        for element in elements_to_render:
                            element_type = element.get('type')
                            element_id = element.get('id')  # Get the ID from the element data
                            new_shape = None
                            try:
                                if element_type in ['text_box', 'text']:
                                    # Pass the element_id to the add_text_box function
                                    new_shape = elements.add_text_box(slide, element, self.style_manager,
                                                                      element_id=element_id)

                                elif element_type == 'image':
                                    if image_keyword := element.get('image_keyword'):
                                        opacity = element.get('style', {}).get('opacity', 1.0)
                                        image_path = image_service.generate_image(image_keyword, opacity)
                                        if image_path:
                                            # Pass the element_id to the add_image function
                                            new_shape = elements.add_image(slide, image_path, element,
                                                                           element_id=element_id)
                                        else:
                                            logging.warning(f"无法为关键词生成图片: '{image_keyword}'。已跳过此元素。")
                                    else:
                                        logging.warning("图片元素缺少 'image_keyword'，已跳过。")

                                elif element_type == 'shape':
                                    # Pass the element_id to the add_shape function
                                    new_shape = elements.add_shape(slide, element, self.style_manager,
                                                                   element_id=element_id)

                                elif element_type == 'chart':
                                    # Pass the element_id to the add_chart function
                                    new_shape = elements.add_chart(slide, element, self.style_manager,
                                                                   element_id=element_id)

                                elif element_type == 'table':
                                    # Pass the element_id to the add_table function
                                    new_shape = elements.add_table(slide, element, self.style_manager,
                                                                   element_id=element_id)

                                # --- 修改2：添加对icon类型的支持 ---
                                elif element_type == 'icon':
                                    # Pass the element_id to the add_icon function
                                    new_shape = elements.add_icon(
                                        slide,
                                        element,
                                        (self.prs.slide_width, self.prs.slide_height),
                                        self.style_manager,
                                        element_id=element_id  # Pass the element_id
                                    )
                                # -----------------------------------
                                else:
                                    logging.warning(f"不支持的元素类型: '{element_type}'。")

                                if element_id and new_shape:
                                    shapes_map[element_id] = new_shape

                            except Exception as e:
                                logging.error(f"渲染类型为 '{element_type}' 的元素失败: {e}", exc_info=True)
                                # --- 第二步：按序列添加动画 ---
                        if 'animation_sequence' in slide_data:
                            for anim_step in slide_data['animation_sequence']:
                                element_id_to_animate = anim_step.get('element_id')
                                anim_data = anim_step.get('animation')

                                if not element_id_to_animate or not anim_data:
                                    continue

                                # 从映射表中查找对应的shape对象
                                shape_to_animate = shapes_map.get(element_id_to_animate)

                                if shape_to_animate:
                                    anim_type = anim_data.get('type')
                                    if anim_type:
                                        # 调用我们强大的动画模块
                                        animation.add_animation(shape_to_animate, anim_type, **anim_data)
                                else:
                                    logging.warning(f"在动画序列中找不到ID为'{element_id_to_animate}'的元素。")

                element_id_to_animate = anim_step.get('element_id')
                anim_data = anim_step.get('animation')

                if not element_id_to_animate or not anim_data:
                    continue

                # 从映射表中查找对应的shape对象
                shape_to_animate = shapes_map.get(element_id_to_animate)

                if shape_to_animate:
                    anim_type = anim_data.get('type')
                    if anim_type:
                        # 调用我们强大的动画模块
                        animation.add_animation(shape_to_animate, anim_type, **anim_data)
                else:
                    logging.warning(f"在动画序列中找不到ID为'{element_id_to_animate}'的元素。")