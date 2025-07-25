# 文件: test_animation_final_modular.py
# [终极模块化版 - 解决了重复创建和扩展性问题]

import logging
import os
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.enum.shapes import MSO_SHAPE

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def get_or_create_main_anim_container(slide_element):
    """
    [模块一：主厨] 获取或创建幻灯片的动画主容器。
    确保 <p:timing> 和主序列只被创建一次。
    :param slide_element: 幻灯片的 lxml 元素。
    :return: (可以附加动画效果的 <p:childTnLst> 节点, 主 <p:timing> 节点)
    """
    timing = slide_element.find(qn('p:timing'))
    if timing is None:
        logging.info("未找到 <p:timing> 节点，正在创建动画骨架...")
        # 使用经过验证的XML骨架，确保结构正确
        timing_xml_str = f"""
        <p:timing {nsdecls('p')}>
            <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
            <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst/>
            </p:cTn></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst><p:bldLst/>
        </p:timing>"""
        timing = parse_xml(timing_xml_str)
        slide_element.append(timing)

    # 返回主序列中用于添加动画的那个 childTnLst 节点，以及timing节点本身
    main_anim_container = timing.find(qn('p:tnLst')).find(qn('p:par')).find(qn('p:cTn')).find(qn('p:childTnLst')).find(
        qn('p:seq')).find(qn('p:cTn')).find(qn('p:childTnLst'))
    return main_anim_container, timing


def build_fly_in_effect(shape_id, next_id, direction, duration_ms):
    """
    [模块二：菜谱] 构建一个独立的“飞入”动画效果XML片段。
    这个函数只关心效果本身，不关心容器。
    :return: 一个代表飞入效果的 <p:par> lxml 元素。
    """
    direction_map = {
        "fromBottom": {"preset": "4", "y_from": "1+#ppt_h/2", "x_from": "#ppt_x"},
        "fromLeft": {"preset": "8", "y_from": "#ppt_y", "x_from": "-#ppt_w"},
        "fromRight": {"preset": "2", "y_from": "#ppt_y", "x_from": "1"},
        "fromTop": {"preset": "1", "y_from": "-#ppt_h", "x_from": "#ppt_x"},
    }
    props = direction_map.get(direction, direction_map["fromBottom"])

    # 完整、精确地构建经过验证的XML结构
    effect_xml_str = """
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:cTn id="{id_3}" fill="hold">
            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
            <p:childTnLst>
                <p:par>
                    <p:cTn id="{id_4}" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="{id_5}" presetID="2" presetClass="entr" presetSubtype="{presetSubtype}" fill="hold" grpId="0" nodeType="clickEffect">
                                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                    <p:childTnLst>
                                        <p:set>
                                            <p:cBhvr>
                                                <p:cTn id="{id_6}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:to><p:strVal val="visible"/></p:to>
                                        </p:set>
                                        <p:anim calcmode="lin" valueType="num">
                                            <p:cBhvr additive="base">
                                                <p:cTn id="{id_7}" dur="{duration}" fill="hold"/>
                                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                                <p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst>
                                                <p:tav tm="0"><p:val><p:strVal val="{anim_x_from}"/></p:val></p:tav>
                                                <p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav>
                                            </p:tavLst>
                                        </p:anim>
                                        <p:anim calcmode="lin" valueType="num">
                                            <p:cBhvr additive="base">
                                                <p:cTn id="{id_8}" dur="{duration}" fill="hold"/>
                                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                                <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst>
                                                <p:tav tm="0"><p:val><p:strVal val="{anim_y_from}"/></p:val></p:tav>
                                                <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
                                            </p:tavLst>
                                        </p:anim>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
    </p:par>
    """.format(
        id_3=str(next_id), id_4=str(next_id + 1), id_5=str(next_id + 2),
        id_6=str(next_id + 3), id_7=str(next_id + 4), id_8=str(next_id + 5),
        spid=shape_id, duration=str(duration_ms),
        presetSubtype=props["preset"], anim_x_from=props["x_from"], anim_y_from=props["y_from"]
    )
    return parse_xml(effect_xml_str)


def build_dissolve_in_effect(shape_id, next_id, duration_ms):
    """
    [新增的菜谱] 构建一个独立的“向内溶解”动画效果XML片段。
    :return: 一个代表向内溶解效果的 <p:par> lxml 元素。
    """
    # 这个XML模板精确复刻了您提供的slide1.xml中spid="5"的动画结构
    effect_xml_str = """
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:cTn id="{id_3}" fill="hold">
            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
            <p:childTnLst>
                <p:par>
                    <p:cTn id="{id_4}" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="{id_5}" presetID="9" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                    <p:childTnLst>
                                        <p:set>
                                            <p:cBhvr>
                                                <p:cTn id="{id_6}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:to><p:strVal val="visible"/></p:to>
                                        </p:set>
                                        <p:animEffect transition="in" filter="dissolve">
                                            <p:cBhvr>
                                                <p:cTn id="{id_7}" dur="{duration}"/>
                                                <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                            </p:cBhvr>
                                        </p:animEffect>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
    </p:par>
    """.format(
        id_3=str(next_id), id_4=str(next_id + 1), id_5=str(next_id + 2),
        id_6=str(next_id + 3), id_7=str(next_id + 4), spid=shape_id, duration=str(duration_ms)
    )
    return parse_xml(effect_xml_str)


# --- 模块三：动画系统的“服务员” (已更新) ---
def add_animation(shape, effect: str, **kwargs):
    """
    为一个形状添加指定的动画效果。这是我们暴露给外部的唯一公共函数。
    """
    logging.info(f"为形状 ID:{shape.shape_id} 添加 '{effect}' 动画...")

    main_anim_container, timing = get_or_create_main_anim_container(shape.part._element)

    # 每个动画效果（一个单击序列）大约需要10个ID的空间
    next_id = len(list(main_anim_container)) * 10 + 3

    effect_fragment = None
    if effect == "flyIn":
        direction = kwargs.get("direction", "fromBottom")
        duration_ms = kwargs.get("duration_ms", 500)
        effect_fragment = build_fly_in_effect(str(shape.shape_id), next_id, direction, duration_ms)
    # --- 新增的扩展点 ---
    elif effect == "dissolveIn":
        duration_ms = kwargs.get("duration_ms", 500)
        effect_fragment = build_dissolve_in_effect(str(shape.shape_id), next_id, duration_ms)
    else:
        logging.warning(f"未知的动画效果类型: '{effect}'")
        return

    if effect_fragment is not None:
        main_anim_container.append(effect_fragment)

        bldLst = timing.find(qn('p:bldLst'))
        if bldLst is None:
            bldLst = OxmlElement('p:bldLst')
            timing.append(bldLst)

        bldP = OxmlElement('p:bldP', {'spid': str(shape.shape_id), 'grpId': '0'})
        bldP.append(OxmlElement('p:bldSub'))
        bldLst.append(bldP)
        logging.info("动画效果已成功注入。")


# --- 测试主程序 (已更新) ---
def run_extended_test():
    """
    测试包含“飞入”和新增“向内溶解”效果的动画系统。
    """
    logging.info("--- 开始扩展版动画功能测试 ---")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "飞入与向内溶解效果"

    # 添加第一个形状：飞入
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1.5))
    shape1.text = "从左侧飞入"
    add_animation(shape1, "flyIn", direction="fromLeft", duration_ms=1000)

    # 添加第二个形状：向内溶解
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(4), Inches(3), Inches(1.5))
    shape2.text = "向内溶解"
    add_animation(shape2, "dissolveIn", duration_ms=500)  # 溶解动画，持续1.5秒

    output_file = "test_animation_dissolve.pptx"
    prs.save(output_file)
    logging.info(f"--- 测试完成 ---")
    logging.info(f"请检查文件 '{os.path.abspath(output_file)}'。")


if __name__ == "__main__":
    run_extended_test()